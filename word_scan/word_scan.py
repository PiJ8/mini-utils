import os
import sys
import re
import argparse

import pypdf
import docx
import pptx
import csv

def argparser(args):
    parser = argparse.ArgumentParser(description='Word count for given directory files')
    parser.add_argument('directory', help='Directory to perform this task')
    parser.add_argument('-r', '--recursive', default=False, action='store_true', help='Recursive for subdirectories')
    # parser.add_argument('-o', '--output', default=False, action='store_true', help='Output result as .csv in current directory')
    return parser.parse_args(args)

# Generated by ChatGPT
def replace_non_alphanumeric(input_string, replacement_char='_'):
    # Use regular expression to replace non-alphanumeric characters with the specified replacement character
    result_string = re.sub(r'[^a-zA-Z0-9 ]', replacement_char, input_string)
    return result_string

def pptx_file_scan(file_path):
    pptx_result = []
    shape_count = 1

    try:
        # Load the PowerPoint presentation
        presentation = pptx.Presentation(file_path)

        # Iterate through slides and text
        for i, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_result = words_count(shape.text, shape_count, file_path)
                    shape_count = shape_count + 1
                    pptx_result = pptx_result + shape_result
    except FileNotFoundError:
        print(f"Error: File '{file_path}' not found.")
    except Exception as ex:
        print(f'Unable to process file: {file_path} with exception {ex}')
    
    return pptx_result

def text_file_scan(file_path):
    text_result = []
    
    try:
        file_handler = open(file_path ,encoding='utf8')
        file_string = file_handler.read()
        text_result = words_count(file_string, 0, file_path)
    except FileNotFoundError:
        print(f"Error: File '{file_path}' not found.")
    except Exception as ex:
        print(f'Unable to process file: {file_path} with exception {ex}')
    
    return text_result

def word_file_scan(file_path):
    word_result = []
    paragraph_count = 1

    # Generated by ChatGPT
    try:
        doc = docx.Document(file_path)
        for paragraph in doc.paragraphs:
            paragraph_string = paragraph.text
            paragraph_result = words_count(paragraph_string, paragraph_count, file_path)
            word_result = word_result + paragraph_result
            paragraph_count = paragraph_count + 1
    except FileNotFoundError:
        print(f"Error: File '{file_path}' not found.")
    except Exception as ex:
        print(f'Unable to process file: {file_path} with exception {ex}')

    return word_result

def pdf_file_scan(file_path):
    pdf_result = []
    page_count = 1

    try:
        reader = pypdf.PdfReader(file_path)
        for page in reader.pages:
            page_string = page.extract_text()
            page_result = words_count(page_string, page_count, file_path)
            pdf_result = pdf_result + page_result
            page_count = page_count + 1
    except FileNotFoundError:
        print(f"Error: File '{file_path}' not found.")
    except Exception as ex:
        print(f'Unable to process file: {file_path} with exception {ex}')
    
    return pdf_result

def words_count(page_string, count=0, file_path=''):
    results = []

    page_string = replace_non_alphanumeric(page_string,' ')
    words = page_string.replace('\n',' ').split(' ')
    words_count_dict = {i: words.count(i) for i in words}
    for key in words_count_dict:
        result_dict = {'file_path': file_path, 
                       'count': count, 
                       # count has different meaning depend on file type
                       # pdf == page | word == paragraph | txt always 0
                       'word': key, 
                       'word_occurence_count': words_count_dict[key]}
        results.append(result_dict)

    return results

def main(arguments):
    
    file_paths = []
    result_in_list = []

    # Scan directory, default current, if --recursive loop subdirectories
    for folderpath, subfolders, filenames in os.walk(arguments.directory):
        # Current directory files
        for filename in filenames:
            file_path = os.path.join(arguments.directory, filename)
            if os.path.isfile(file_path):
                file_paths.append(file_path)

        # Subdirectories files
        if arguments.recursive:
            for filename in os.listdir(folderpath):
                file_path = os.path.join(folderpath, filename)
                if os.path.isfile(file_path):
                    file_paths.append(file_path)
    
    for file_path in file_paths:
        _, file_extension = os.path.splitext(file_path)

        if file_extension == '.txt':
            result= text_file_scan(file_path)
        elif file_extension == '.docx':
            result= word_file_scan(file_path)
        elif file_extension == '.pdf':
            result= pdf_file_scan(file_path)
        elif file_extension == '.pptx':
            result= pptx_file_scan(file_path)
        else:
            print(f"Unsupported file type {file_extension} for {file_path}")
        
        result_in_list = result_in_list + result

    # search_dict = ['something', 'manything']

    csv_file_name = 'word_scan_result.tmp.csv'
    
    try:
        # Open the CSV file in write mode
        with open(csv_file_name, 'w', newline='') as csvfile:
            # Create a CSV writer object
            csv_writer = csv.writer(csvfile)

            # Write the header (keys of the dictionary)
            csv_writer.writerow(result_in_list[0].keys())

            for item in result_in_list:
                # Write the values of the dictionary
                csv_writer.writerow(item.values())
            
            print(f'Data has been written to {csv_file_name}.')
    except Exception as e:
        print(f"An error occurred: {e}")

    # try:
    #     with open('testresult.tmp.txt', 'w') as file:
    #         for item in result_in_list:
    #             file.write(str(item) + '\n')
    #     print(f"List saved to testresult successfully.")
    # except Exception as e:
    #     print(f"An error occurred: {e}")



if __name__ == '__main__':
    argv = argparser(sys.argv[1:])
    main(argv)